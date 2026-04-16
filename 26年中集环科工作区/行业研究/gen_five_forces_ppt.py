# -*- coding: utf-8 -*-
"""
生成罐箱产业波特五力分析PPT（1页）
用于2026年Q1总裁工作报告
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ── 配色方案 ──────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1B, 0x3A, 0x5C)   # 深蓝（标题/边框）
MID_BLUE    = RGBColor(0x2C, 0x5F, 0x8A)   # 中蓝（中心）
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)   # 淡蓝（卡片背景）
RED_ACCENT  = RGBColor(0xC0, 0x39, 0x2B)   # 红色（高威胁）
AMBER       = RGBColor(0xE6, 0x7E, 0x22)   # 橙色（中高威胁）
GREEN       = RGBColor(0x27, 0xAE, 0x60)   # 绿色（低威胁）
DARK_GREEN  = RGBColor(0x1E, 0x8A, 0x4C)   # 深绿
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x2C, 0x2C, 0x2C)
GRAY        = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY  = RGBColor(0xEC, 0xF0, 0xF1)
ARROW_GRAY  = RGBColor(0x95, 0xA5, 0xA6)

# ── 评级颜色映射 ─────────────────────────────────────
RATING_COLORS = {
    "高":   RED_ACCENT,
    "中高": AMBER,
    "中":   AMBER,
    "中低": GREEN,
    "低":   GREEN,
    "低-中": GREEN,
}

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1.5)):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # 圆角调节 (0~100000, 越小越圆)
    shape.adjustments[0] = 0.08
    return shape

def set_text(tf, text, font_size=9, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    """设置文本框内容"""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Microsoft YaHei"
    return p

def add_paragraph(tf, text, font_size=8, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, space_before=Pt(2)):
    """追加段落"""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Microsoft YaHei"
    return p

def add_force_card(slide, left, top, width, height,
                   title, rating, rating_color, bullets, border_color):
    """添加一个力量分析卡片"""
    # 背景卡片
    card = add_rounded_rect(slide, left, top, width, height, WHITE, border_color, Pt(2))
    tf = card.text_frame
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(4)
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.word_wrap = True

    # 标题行
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(1)
    run = p.add_run()
    run.text = title
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    run.font.name = "Microsoft YaHei"

    # 评级标签
    run2 = p.add_run()
    run2.text = f"  【{rating}】"
    run2.font.size = Pt(10)
    run2.font.bold = True
    run2.font.color.rgb = rating_color
    run2.font.name = "Microsoft YaHei"

    # 要点
    for bullet in bullets:
        add_paragraph(tf, f"• {bullet}", font_size=7.5, color=BLACK, space_before=Pt(2))

    return card

def add_arrow(slide, start_left, start_top, end_left, end_top, color=ARROW_GRAY):
    """添加带箭头的连接线"""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_left, start_top,
        end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    # 箭头终点
    connector.end_x = end_left
    connector.end_y = end_top
    return connector

def main():
    prs = Presentation()
    # 16:9 宽屏
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # ── 背景色 ──
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)

    # ── 标题栏 ──
    title_bar = add_rounded_rect(
        slide, Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.6),
        DARK_BLUE, None
    )
    tf = title_bar.text_frame
    tf.margin_left = Pt(12)
    set_text(tf, "ISO罐式集装箱产业 — 波特五力分析（Porter's Five Forces）",
             font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
    # 副标题
    add_paragraph(tf, "2026年Q1经营数据视角  |  全球罐箱保有量 899,044台  |  新造产能 ~12万台/年  |  行业产能利用率 24%",
                  font_size=9, color=RGBColor(0xB0, 0xC4, 0xDE), alignment=PP_ALIGN.LEFT, space_before=Pt(1))

    # ══════════════════════════════════════════════════════
    # 布局参数
    # ══════════════════════════════════════════════════════
    card_w = Inches(3.6)
    card_h = Inches(2.15)
    center_w = Inches(3.8)
    center_h = Inches(2.5)

    # 中心点
    cx = Inches(6.666) - center_w / 2   # 水平居中
    cy = Inches(3.85) - center_h / 2    # 垂直略偏上

    # 上方（新进入者）
    top_x = Inches(6.666) - card_w / 2
    top_y = Inches(1.0)

    # 下方（替代品）
    bot_x = Inches(6.666) - card_w / 2
    bot_y = Inches(6.15)

    # 左侧（供应商）
    left_x = Inches(0.4)
    left_y = Inches(3.85) - card_h / 2

    # 右侧（客户）
    right_x = Inches(13.333) - card_w - Inches(0.4)
    right_y = Inches(3.85) - card_h / 2

    # ── 1. 新进入者威胁（上方）──  低 ──
    add_force_card(
        slide, top_x, top_y, card_w, card_h,
        title="⬆ 新进入者威胁",
        rating="低",
        rating_color=GREEN,
        bullets=[
            "资质壁垒：UN/IMDG/ADR/TSG多国认证，周期18-24月",
            "资本壁垒：特罐单条产线投资9,980万元",
            "规模经济：行业产能利用率仅24%，新进入者难盈利",
            "技术积累：焊接工艺经验20+年，专有设计难复制",
            "Q1信号：全球仅新造28,521台（15年低位）",
        ],
        border_color=GREEN,
    )

    # ── 2. 替代品威胁（下方）──  低-中 ──
    add_force_card(
        slide, bot_x, bot_y, card_w, card_h - Inches(0.15),
        title="⬇ 替代品威胁",
        rating="低-中",
        rating_color=GREEN,
        bullets=[
            "公路罐车：仅适用短途，不可多式联运",
            "铁路罐车：大宗散货专用，路径固定",
            "Flexitank：仅非危险品液体，一次性使用",
            "ISO罐箱优势：多式联运+可重复使用+全球标准化",
            "Q1信号：全球保有量同比+1.93%，化工贸易持续增长",
        ],
        border_color=GREEN,
    )

    # ── 3. 供应商议价能力（左侧）──  中低 ──
    add_force_card(
        slide, left_x, left_y, card_w, card_h,
        title="⬅ 供应商议价能力",
        rating="中低",
        rating_color=GREEN,
        bullets=[
            "316L不锈钢（占成本40-50%）：中国产能过剩",
            "  · 国内5,000万吨产能 vs 3,250万吨需求",
            "  · 国内 $1.70/kg vs 国际 $4.60-5.10/kg",
            "中集为全球最大单一316L采购商，议价力强",
            "阀门：Fort Vale/Perolo双寡头，但中集自制率↑",
            "  · 人孔凸缘自制80%、V型防波板100%",
        ],
        border_color=DARK_GREEN,
    )

    # ── 4. 客户议价能力（右侧）──  中高 ──
    add_force_card(
        slide, right_x, right_y, card_w, card_h,
        title="➡ 客户议价能力",
        rating="中高",
        rating_color=AMBER,
        bullets=[
            "租赁商（占保有量43%）：~15%闲置率，扩张放缓",
            "运营商（占使用量70%）：Hoyer/Stolt集中采购",
            "Q1: 运营商+租赁商签单占比88.7%",
            "标准罐均价 USD 1.31-1.38万（持续承压）",
            "订单+204% 但收入仅+120% → 以量换价明显",
            "对策：向特种罐/智能罐/后市场差异化转型",
        ],
        border_color=AMBER,
    )

    # ── 5. 行业内竞争（中心）──  高 ──
    center_card = add_rounded_rect(
        slide, cx, cy, center_w, center_h,
        MID_BLUE, DARK_BLUE, Pt(2.5)
    )
    tf = center_card.text_frame
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(6)
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.word_wrap = True

    # 中心标题
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = "行业内竞争"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Microsoft YaHei"
    run2 = p.add_run()
    run2.text = "  【高】"
    run2.font.size = Pt(14)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(0xFF, 0x6B, 0x6B)
    run2.font.name = "Microsoft YaHei"

    center_bullets = [
        "Top 7占全球产量93%，中集市场份额~50%",
        "产能利用率仅24%（中集Q1: 38%）→ 价格战激烈",
        "Q1标准罐段亏损 -1,581万，毛利承压",
        "NT Tank/JJAP低价策略侵蚀份额",
        "分化路径：特种罐毛利16% vs 标准罐微利",
        "Q1订单爆发+204%，3月单月扭亏 → 周期回暖信号",
    ]
    for b in center_bullets:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(7.5)
        run.font.color.rgb = RGBColor(0xE8, 0xF0, 0xF8)
        run.font.name = "Microsoft YaHei"

    # ── 箭头（用三角形模拟方向指示）──
    arrow_size = Inches(0.3)

    # 上 → 中心（向下箭头）
    slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(6.666) - arrow_size/2, top_y + card_h + Inches(0.05),
        arrow_size, Inches(0.35)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ARROW_GRAY
    slide.shapes[-1].line.fill.background()

    # 下 → 中心（向上箭头）
    slide.shapes.add_shape(
        MSO_SHAPE.UP_ARROW,
        Inches(6.666) - arrow_size/2, bot_y - Inches(0.4),
        arrow_size, Inches(0.35)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ARROW_GRAY
    slide.shapes[-1].line.fill.background()

    # 左 → 中心（向右箭头）
    slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        left_x + card_w + Inches(0.05), Inches(3.85) - arrow_size/2,
        Inches(0.45), arrow_size
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ARROW_GRAY
    slide.shapes[-1].line.fill.background()

    # 右 → 中心（向左箭头）
    slide.shapes.add_shape(
        MSO_SHAPE.LEFT_ARROW,
        right_x - Inches(0.5), Inches(3.85) - arrow_size/2,
        Inches(0.45), arrow_size
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ARROW_GRAY
    slide.shapes[-1].line.fill.background()

    # ── 底部综合评估栏 ──
    footer = add_rounded_rect(
        slide, Inches(0.3), Inches(6.85), Inches(12.733), Inches(0.5),
        DARK_BLUE, None
    )
    tf = footer.text_frame
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "综合评估："
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Microsoft YaHei"
    run2 = p.add_run()
    run2.text = ("行业竞争烈度高，客户议价能力强，但进入壁垒与替代威胁低→行业格局有利于在位龙头。"
                 "中集环科核心战略：① 特种罐+智能罐差异化提升定价权  ② 后市场服务构建客户粘性  "
                 "③ 316L成本优势筑底护城河  ④ 产能弹性应对周期波动（Q1产能利用率38%→目标55%+）")
    run2.font.size = Pt(8)
    run2.font.color.rgb = RGBColor(0xD6, 0xE8, 0xF7)
    run2.font.name = "Microsoft YaHei"

    # ── 右下角数据来源标注 ──
    src_box = slide.shapes.add_textbox(
        Inches(10.5), Inches(7.15), Inches(2.5), Inches(0.25)
    )
    tf = src_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "数据来源：Q1经营指标 / ITCO / 行业研究报告"
    run.font.size = Pt(6)
    run.font.color.rgb = GRAY
    run.font.name = "Microsoft YaHei"

    # ── 保存 ──
    output = Path(__file__).parent / "罐箱产业五力分析_Q1_2026.pptx"
    prs.save(str(output))
    print(f"PPT saved to: {output}")
    print(f"Slide count: {len(prs.slides)}")

if __name__ == "__main__":
    main()
