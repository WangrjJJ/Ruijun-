#!/usr/bin/env python3
"""
JARVIS Ingestion Agent — 文件摄入代理
扫描文件 → 文本提取 → 哈希去重 → 输出 ingestion_events.jsonl

用法:
  python3 ingestion_agent.py                    # 扫描所有已知源
  python3 ingestion_agent.py --file <path>       # 单文件摄入
  python3 ingestion_agent.py --dir <path>        # 扫描目录
  python3 ingestion_agent.py --scan              # 全量扫描（忽略去重）
  python3 ingestion_agent.py --stats             # 查看摄入统计
"""

import os
import sys
import json
import hashlib
import argparse
from datetime import datetime, timedelta
from pathlib import Path

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.join(SCRIPT_DIR, "data")
os.makedirs(DATA_DIR, exist_ok=True)

EVENTS_FILE = os.path.join(DATA_DIR, "ingestion_events.jsonl")
PROCESSED_FILE = os.path.join(DATA_DIR, "processed_files.json")

# ── 扫描路径配置 ──────────────────────────────────────────────────
# 默认扫描的目录及其标签
SCAN_PATHS = [
    # Obsidian Vault 26年工作区 (markdown notes)
    {
        "path": r"C:\Users\01455310\Documents\Obsidian Vault\26年中集环科工作区",
        "label": "26年工作区",
        "recursive": True,
    },
    # Desktop 26年工作文件 — 根目录独立文件 (PPTX/XLSX/DOCX)
    {
        "path": r"C:\Users\01455310\Desktop\26年工作文件",
        "label": "工作文件根目录",
        "recursive": False,
    },
    # Desktop 子目录
    {
        "path": r"C:\Users\01455310\Desktop\26年工作文件\26年季度经营分析会",
        "label": "季度经营分析",
        "recursive": True,
    },
    {
        "path": r"C:\Users\01455310\Desktop\26年工作文件\26年总裁工作报告",
        "label": "总裁报告",
        "recursive": True,
    },
    {
        "path": r"C:\Users\01455310\Desktop\26年工作文件\特罐搬迁项目",
        "label": "特罐搬迁",
        "recursive": True,
    },
    {
        "path": r"C:\Users\01455310\Desktop\26年工作文件\绩效考核",
        "label": "绩效考核",
        "recursive": True,
    },
    {
        "path": r"C:\Users\01455310\Desktop\26年工作文件\_inbox",
        "label": "收件箱",
        "recursive": True,
    },
]

# 支持的文件扩展名及解析器映射
SUPPORTED_EXT = {
    ".md": "markdown",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    ".docx": "docx",
    ".doc": "doc",
    ".pdf": "pdf",
    ".txt": "txt",
    ".csv": "csv",
    ".py": "python",
}


def file_fingerprint(filepath: str) -> str:
    """基于 mtime + size 的哈希，用于去重"""
    try:
        stat = os.stat(filepath)
        raw = f"{stat.st_mtime:.6f}|{stat.st_size}"
        return hashlib.sha256(raw.encode()).hexdigest()[:16]
    except OSError:
        return ""


def load_processed() -> dict:
    """加载已处理文件注册表: {fingerprint: {path, time, label}}"""
    if os.path.exists(PROCESSED_FILE):
        with open(PROCESSED_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def save_processed(registry: dict):
    with open(PROCESSED_FILE, "w", encoding="utf-8") as f:
        json.dump(registry, f, ensure_ascii=False, indent=2)


def extract_text_markdown(filepath: str) -> str:
    """提取 markdown 文件全文"""
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(filepath, "r", encoding="gbk") as f:
            return f.read()


def extract_text_pptx(filepath: str) -> str:
    """深度提取 PPTX 文本 + 表格结构 + 备注"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation(filepath)
        parts = []

        # 元数据
        props = prs.core_properties
        meta_parts = []
        if props.title:
            meta_parts.append(f"标题: {props.title}")
        if props.author:
            meta_parts.append(f"作者: {props.author}")
        if props.last_modified_by:
            meta_parts.append(f"最后修改: {props.last_modified_by}")
        if props.subject:
            meta_parts.append(f"主题: {props.subject}")
        if meta_parts:
            parts.append("## 文档信息\n" + "\n".join(meta_parts))

        # 全局指标收集
        all_metrics = []
        all_tables = []

        for i, slide in enumerate(prs.slides, 1):
            slide_num = 0
            # 跳过隐藏幻灯片
            try:
                if slide.slide_layout and "hidden" in slide.slide_layout.name.lower():
                    continue
            except Exception:
                pass

            slide_text = []
            slide_metrics = []

            for shape in slide.shapes:
                # 文本提取
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        slide_text.append(text)
                        # 识别关键指标: 含数字+单位/百分比
                        if any(kw in text for kw in ["%", "亿", "万", "同比", "环比",
                                                       "完成率", "渗透率", "IRR", "NPV",
                                                       "收入", "利润", "毛利", "产能",
                                                       "人效", "KPI", "偏差"]):
                            if any(c.isdigit() for c in text):
                                slide_metrics.append({
                                    "slide": i,
                                    "text": text,
                                    "context": slide_text[-2] if len(slide_text) > 1 else "",
                                })
                                all_metrics.append(slide_metrics[-1])

                # 表格深度提取
                if shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_cells)

                    if table_data:
                        all_tables.append({
                            "slide": i,
                            "rows": len(table_data),
                            "cols": len(table_data[0]) if table_data else 0,
                            "data": table_data,
                        })
                        # 渲染为 markdown 表格
                        header = table_data[0]
                        sep = ["---"] * len(header)
                        md_table = [f"| {' | '.join(header)} |"]
                        md_table.append(f"| {' | '.join(sep)} |")
                        for row in table_data[1:]:
                            padded = row + [""] * (len(header) - len(row))
                            md_table.append(f"| {' | '.join(padded[:len(header)])} |")
                        slide_text.append("\n" + "\n".join(md_table[:15]))  # 最多15行

                # 图表数据（通过 chart 对象）
                if shape.has_chart:
                    try:
                        chart = shape.chart
                        slide_text.append(f"[图表: {chart.chart_type} — {chart.has_legend and '含图例' or ''}]")
                    except Exception:
                        slide_text.append("[图表]")

            # 提取备注
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text.append(f"[备注] {notes}")

            if slide_text:
                header = f"## Slide {i}"
                if slide_metrics:
                    header += f" (含{len(slide_metrics)}项指标)"
                parts.append(header + "\n" + "\n".join(slide_text[:50]))  # 每slide最多50行

        # 汇总
        summary_parts = []
        if all_metrics:
            summary_parts.append(f"\n## 关键指标汇总 (共{len(all_metrics)}项)")
            for m in all_metrics[:15]:
                summary_parts.append(f"- [Slide {m['slide']}] {m['text']}")

        if all_tables:
            summary_parts.append(f"\n## 表格汇总 (共{len(all_tables)}个)")
            for t in all_tables:
                summary_parts.append(f"- Slide {t['slide']}: {t['rows']}行 × {t['cols']}列")
                if t["data"]:
                    summary_parts.append(f"  表头: {' | '.join(t['data'][0][:5])}")

        return "\n\n".join(parts + summary_parts)
    except ImportError:
        return "[PPTX提取需要python-pptx]"
    except Exception as e:
        return f"[PPTX提取失败: {e}]"


def extract_text_xlsx(filepath: str) -> str:
    """深度提取 XLSX — 识别关键指标、表格结构、数据趋势"""
    try:
        import openpyxl
        from openpyxl.utils import get_column_letter

        wb = openpyxl.load_workbook(filepath, data_only=True)
        parts = []

        # 文档属性
        props = wb.properties
        meta_parts = []
        if props.title:
            meta_parts.append(f"标题: {props.title}")
        if props.creator:
            meta_parts.append(f"创建者: {props.creator}")
        if props.lastModifiedBy:
            meta_parts.append(f"最后修改: {props.lastModifiedBy}")
        if meta_parts:
            parts.append("## 文档信息\n" + "\n".join(meta_parts))

        all_metrics = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            if ws.max_row == 1 and ws.max_column == 1 and ws["A1"].value is None:
                continue  # 跳过空 sheet

            sheet_parts = [f"## Sheet: {sheet_name} ({ws.max_row}行 × {ws.max_column}列)"]

            # 提取表头行（前5行中找有内容的）
            header_row = 0
            for r in range(1, min(6, ws.max_row + 1)):
                filled = sum(1 for c in range(1, ws.max_column + 1)
                            if ws.cell(row=r, column=c).value is not None)
                if filled >= 2:
                    header_row = r
                    break

            # 提取表格结构
            if header_row:
                headers = []
                for c in range(1, ws.max_column + 1):
                    val = ws.cell(row=header_row, column=c).value
                    headers.append(str(val)[:30] if val is not None else "")
                sheet_parts.append(f"表头: {' | '.join(headers[:8])}")

            # 扫描指标行
            metric_rows = []
            for r in range(1, ws.max_row + 1):
                row_values = []
                for c in range(1, ws.max_column + 1):
                    val = ws.cell(row=r, column=c).value
                    if val is not None:
                        row_values.append(str(val))
                row_str = " ".join(row_values)

                # 识别含关键指标的行
                is_metric = any(kw in row_str for kw in [
                    "%", "亿", "万", "同比", "环比", "完成率", "渗透率",
                    "收入", "利润", "毛利", "产能", "人效", "KPI",
                    "偏差", "预算", "实际", "目标", "同比增长", "环比增长",
                    "营业收入", "净利润", "人均产值", "毛利率", "净利率",
                ])
                if is_metric and any(c.isdigit() for c in row_str):
                    cells = [str(ws.cell(row=r, column=c).value or "")
                            for c in range(1, min(ws.max_column + 1, 9))]
                    metric_rows.append({
                        "row": r,
                        "cells": " | ".join(cells),
                        "sheet": sheet_name,
                    })
                    all_metrics.append(metric_rows[-1])

            # 输出表头 + 指标行
            if metric_rows:
                sheet_parts.append(f"\n关键指标行 ({len(metric_rows)}项):")
                for m in metric_rows[:20]:
                    sheet_parts.append(f"  Row {m['row']}: {m['cells'][:120]}")
            else:
                # 无指标匹配时输出前20行数据
                sample_rows = []
                for r in range(1, min(21, ws.max_row + 1)):
                    cell_vals = [str(ws.cell(row=r, column=c).value or "")
                                for c in range(1, min(ws.max_column + 1, 6))]
                    row_str = " | ".join(cell_vals)
                    if row_str.replace(" | ", "").strip():
                        sample_rows.append(f"  Row {r}: {row_str}")
                if sample_rows:
                    sheet_parts.append("\n数据预览:")
                    sheet_parts.extend(sample_rows[:10])

            parts.append("\n".join(sheet_parts))

        wb.close()

        # 全局指标汇总
        if all_metrics:
            parts.append(f"\n\n## 指标汇总 ({len(all_metrics)}项)")
            for m in all_metrics[:20]:
                parts.append(f"- [{m['sheet']}] Row {m['row']}: {m['cells'][:100]}")

        return "\n\n".join(parts)
    except ImportError:
        return "[XLSX提取需要openpyxl]"
    except Exception as e:
        return f"[XLSX提取失败: {e}]"


def extract_text_docx(filepath: str) -> str:
    """提取 DOCX 文本"""
    try:
        from docx import Document
        doc = Document(filepath)
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        # 提取表格
        for i, table in enumerate(doc.tables):
            table_parts = []
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells
                )
                if row_text.strip():
                    table_parts.append(row_text)
            if table_parts:
                parts.append(f"\n[Table {i+1}]\n" + "\n".join(table_parts))
        return "\n\n".join(parts)
    except Exception as e:
        return f"[DOCX提取失败: {e}]"


def extract_text_pdf(filepath: str) -> str:
    """提取 PDF 文本（简易版，使用 PyPDF2）"""
    try:
        import PyPDF2
        text_parts = []
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for i, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    text_parts.append(text)
                if i > 50:
                    text_parts.append("... (truncated)")
                    break
        return "\n\n".join(text_parts)
    except ImportError:
        return "[PDF提取需要安装PyPDF2: pip install PyPDF2]"
    except Exception as e:
        return f"[PDF提取失败: {e}]"


def extract_text_plain(filepath: str) -> str:
    """提取纯文本文件"""
    for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
        try:
            with open(filepath, "r", encoding=encoding) as f:
                return f.read()[:10000]
        except UnicodeDecodeError:
            continue
    return "[无法解码文件]"


def extract_text_doc(filepath: str) -> str:
    """提取 .doc (旧版Word) 文本 — 通过 win32com 或 python-docx 降级"""
    # 先尝试 python-docx（某些 .doc 实际上是 docx 格式）
    try:
        from docx import Document
        Document(filepath)
        # 如果成功打开, 恢复到 docx 解析
        parts = []
        doc = Document(filepath)
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        return "\n\n".join(parts)
    except Exception:
        pass

    # 使用 win32com 通过 Word 转换
    try:
        import win32com.client
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(filepath, ReadOnly=True)
        text = doc.Content.Text
        doc.Close()
        word.Quit()
        return text[:20000] if text else ""
    except ImportError:
        return "[.doc提取需要: pip install pywin32 并安装 Microsoft Word]"
    except Exception as e:
        return f"[.doc提取失败: {e}]"


EXTRACTORS = {
    ".md": extract_text_markdown,
    ".pptx": extract_text_pptx,
    ".xlsx": extract_text_xlsx,
    ".docx": extract_text_docx,
    ".doc": extract_text_doc,
    ".pdf": extract_text_pdf,
    ".txt": extract_text_plain,
    ".csv": extract_text_plain,
    ".py": extract_text_plain,
}


def generate_summary(text: str, max_len: int = 200, file_type: str = "") -> str:
    """从提取文本生成简短摘要"""
    # 对 markdown 文件跳过 frontmatter
    content = text
    if file_type == "markdown" and text.startswith("---"):
        end = text.find("---", 3)
        if end > 0:
            content = text[end + 3:].strip()
    lines = [l.strip() for l in content.split("\n") if l.strip() and not l.startswith("[")]
    summary = " ".join(lines[:5])
    if len(summary) > max_len:
        summary = summary[:max_len] + "..."
    return summary


def extract_topics(text: str, max_topics: int = 5) -> list:
    """从文本中提取可能的关键话题词"""
    # 基于 frontmatter tags 和关键词出现频率的简易提取
    topics = []
    # 尝试提取 frontmatter tags
    if text.startswith("---"):
        end = text.find("---", 3)
        if end > 0:
            fm = text[3:end].strip()
            for line in fm.split("\n"):
                if line.startswith("tags:") or line.startswith("tags:"):
                    continue
                if "tags:" in fm:
                    tag_section = fm.split("tags:")[1].split("\n")[0]
                    for tag in tag_section.replace("[", "").replace("]", "").split(","):
                        t = tag.strip().strip("'").strip('"')
                        if t and len(t) > 1:
                            topics.append(t)
                    break
    return topics[:max_topics]


def ingest_file(filepath: str, label: str = "", force: bool = False) -> dict | None:
    """摄入单个文件，返回事件或 None（如果已处理）"""
    if not os.path.isfile(filepath):
        return None

    ext = Path(filepath).suffix.lower()
    if ext not in SUPPORTED_EXT:
        return None

    fp = file_fingerprint(filepath)
    if not fp:
        return None

    registry = load_processed()
    if not force and fp in registry:
        return None

    extractor = EXTRACTORS.get(ext)
    if not extractor:
        return None

    extracted = extractor(filepath)
    if not extracted or extracted.startswith("["):
        return None

    summary = generate_summary(extracted, file_type=SUPPORTED_EXT[ext])
    topics = extract_topics(extracted)

    event = {
        "timestamp": datetime.now().strftime("%Y-%m-%dT%H:%M:%S"),
        "fingerprint": fp,
        "source_path": filepath,
        "source_label": label,
        "file_type": SUPPORTED_EXT[ext],
        "file_name": os.path.basename(filepath),
        "summary": summary,
        "key_topics": topics,
        "extracted_length": len(extracted),
        "extracted_text": extracted,
    }

    # 写 JSONL
    with open(EVENTS_FILE, "a", encoding="utf-8") as f:
        f.write(json.dumps(event, ensure_ascii=False) + "\n")

    # 更新注册表
    registry[fp] = {
        "path": filepath,
        "label": label,
        "time": event["timestamp"],
        "file_name": event["file_name"],
    }
    save_processed(registry)

    return event


def scan_directory(dir_path: str, label: str = "", recursive: bool = True) -> list:
    """扫描目录下所有支持的文件，返回摄入事件列表"""
    if not os.path.isdir(dir_path):
        return []
    events = []
    pattern = "*" if not recursive else "**/*"
    for ext in SUPPORTED_EXT:
        for filepath in Path(dir_path).glob(f"{pattern}{ext}"):
            try:
                event = ingest_file(str(filepath), label)
                if event:
                    events.append(event)
            except Exception as e:
                print(f"  [WARN] {filepath}: {e}", file=sys.stderr)
    return events


def cmd_scan(args=None):
    """全量扫描所有已配置路径"""
    all_events = []
    for cfg in SCAN_PATHS:
        path = cfg["path"]
        if not os.path.exists(path):
            print(f"  [SKIP] 不存在: {path}", file=sys.stderr)
            continue
        print(f"[SCAN] {cfg['label']}: {path}", file=sys.stderr)
        events = scan_directory(path, cfg["label"], cfg.get("recursive", True))
        all_events.extend(events)
        print(f"  → {len(events)} 新文件", file=sys.stderr)
    return all_events


def cmd_file(filepath: str, label: str = "手动摄入", force: bool = False):
    """单文件摄入"""
    if not os.path.exists(filepath):
        print(f"文件不存在: {filepath}", file=sys.stderr)
        return None
    event = ingest_file(filepath, label, force=force)
    if event:
        print(f"[INGEST] {event['file_name']} → {event['summary'][:80]}", file=sys.stderr)
    else:
        print(f"[SKIP] 已处理或无变化: {filepath}", file=sys.stderr)
    return event


def cmd_stats():
    """摄入统计"""
    if not os.path.exists(EVENTS_FILE):
        print("暂无摄入事件")
        return

    events = []
    with open(EVENTS_FILE, "r", encoding="utf-8") as f:
        for line in f:
            if line.strip():
                events.append(json.loads(line))

    registry = load_processed()
    print(f"摄入统计:")
    print(f"  总摄入事件: {len(events)}")
    print(f"  已处理文件数: {len(registry)}")

    # 按类型统计
    by_type = {}
    for e in events:
        t = e.get("file_type", "unknown")
        by_type[t] = by_type.get(t, 0) + 1
    print(f"  按文件类型: {by_type}")

    # 按来源统计
    by_label = {}
    for e in events:
        lbl = e.get("source_label", "unknown")
        by_label[lbl] = by_label.get(lbl, 0) + 1
    print(f"  按来源: {by_label}")

    # 最近24小时
    cutoff = (datetime.now() - timedelta(hours=24)).strftime("%Y-%m-%dT%H:%M:%S")
    recent = [e for e in events if e["timestamp"] >= cutoff]
    print(f"  最近24小时: {len(recent)} 事件")
    for e in recent:
        print(f"    - [{e['file_type']}] {e['file_name']}: {e['summary'][:60]}")


def cmd_recent(hours: int = 24) -> list:
    """获取最近N小时的摄入事件"""
    if not os.path.exists(EVENTS_FILE):
        return []
    cutoff = (datetime.now() - timedelta(hours=hours)).strftime("%Y-%m-%dT%H:%M:%S")
    events = []
    with open(EVENTS_FILE, "r", encoding="utf-8") as f:
        for line in f:
            if line.strip():
                event = json.loads(line)
                if event["timestamp"] >= cutoff:
                    events.append(event)
    return events


def cmd_watch(interval: int = 300):
    """持续监听模式 — 每 N 秒扫描一次已知目录的新文件"""
    import time
    print(f"JARVIS 文件监听已启动 · 间隔 {interval}s · Ctrl+C 停止", file=sys.stderr)
    print(f"监听目录:", file=sys.stderr)
    for cfg in SCAN_PATHS:
        exists = os.path.exists(cfg["path"])
        status = "✓" if exists else "✗ (不存在)"
        print(f"  {status} {cfg['path']}", file=sys.stderr)
    print("", file=sys.stderr)

    iteration = 0
    try:
        while True:
            iteration += 1
            ts = datetime.now().strftime("%H:%M:%S")
            print(f"[{ts}] 扫描 # {iteration}...", file=sys.stderr)
            events = cmd_scan()
            if events:
                print(f"  → {len(events)} 个新文件摄入", file=sys.stderr)
                for e in events:
                    print(f"    [{e['file_type']}] {e['file_name']}", file=sys.stderr)
            else:
                print(f"  → 无新文件", file=sys.stderr)
            time.sleep(interval)
    except KeyboardInterrupt:
        print(f"\n监听已停止 · 共扫描 {iteration} 次", file=sys.stderr)


def main():
    parser = argparse.ArgumentParser(description="JARVIS 摄入代理")
    group = parser.add_mutually_exclusive_group()
    group.add_argument("--file", help="摄入单个文件")
    group.add_argument("--dir", help="扫描目录")
    group.add_argument("--scan", action="store_true", help="全量扫描所有已知路径")
    group.add_argument("--stats", action="store_true", help="显示摄入统计")
    group.add_argument("--recent", type=int, metavar="HOURS", help="获取最近N小时事件(JSON输出)")
    group.add_argument("--watch", action="store_true", help="持续监听模式")
    parser.add_argument("--force", action="store_true", help="强制重新摄入（忽略去重）")
    parser.add_argument("--label", default="", help="来源标签")
    parser.add_argument("--interval", type=int, default=300,
                       help="监听间隔秒数（默认300）")
    args = parser.parse_args()

    if args.stats:
        cmd_stats()
    elif args.recent:
        events = cmd_recent(args.recent)
        print(json.dumps(events, ensure_ascii=False, indent=2))
    elif args.file:
        cmd_file(args.file, args.label or "手动摄入", args.force)
    elif args.dir:
        events = scan_directory(args.dir, args.label or os.path.basename(args.dir))
        print(f"摄入完成: {len(events)} 个新文件", file=sys.stderr)
        for e in events:
            print(f"  [{e['file_type']}] {e['file_name']}", file=sys.stderr)
    elif args.watch:
        cmd_watch(args.interval)
    elif args.scan:
        events = cmd_scan()
        print(f"\n总计摄入: {len(events)} 个新文件", file=sys.stderr)
    else:
        # 默认：快速扫描（检查已知路径的新文件）
        print("快速扫描模式...", file=sys.stderr)
        events = cmd_scan()
        print(f"完成: {len(events)} 个新文件", file=sys.stderr)


if __name__ == "__main__":
    main()
